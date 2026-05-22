"""Generate NoQx_Launch_Plan.pptx — operator-shareable launch deck for the client.

Run from repo root:
  python3 scripts/generate-launch-ppt.py

Output: NoQx_Launch_Plan.pptx in repo root.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "NoQx_Launch_Plan.pptx"
LOGO = ROOT / "public" / "icons" / "icon-512.png"

# NoQx brand
PURPLE = RGBColor(0x7c, 0x3a, 0xed)
PURPLE_DARK = RGBColor(0x1a, 0x15, 0x30)
INK_DARK = RGBColor(0x1e, 0x29, 0x3b)
INK_2 = RGBColor(0x47, 0x55, 0x69)
INK_3 = RGBColor(0x64, 0x74, 0x8b)
WHITE = RGBColor(0xff, 0xff, 0xff)
GREEN = RGBColor(0x16, 0xa3, 0x4a)
AMBER = RGBColor(0xd9, 0x77, 0x06)
RED = RGBColor(0xdc, 0x26, 0x26)
LIGHT_BG = RGBColor(0xf5, 0xf3, 0xff)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]

# ── helpers ──────────────────────────────────────────────────────────────
def add_solid_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg

def add_text(slide, left, top, width, height, text, *, size=18, bold=False,
             color=INK_DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb

def add_pill(slide, left, top, width, height, text, fill, text_color=WHITE, size=11, bold=True):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill
    pill.line.fill.background()
    tf = pill.text_frame
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = text_color
    return pill

def add_header_strip(slide, title, subtitle=None):
    # Top purple strip
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    strip.fill.solid(); strip.fill.fore_color.rgb = PURPLE
    strip.line.fill.background()
    add_text(slide, Inches(0.5), Inches(0.18), Inches(11), Inches(0.5),
             title, size=24, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.6), Inches(11), Inches(0.3),
                 subtitle, size=12, color=WHITE)
    # Page-number / brand on right
    add_text(slide, Inches(11.8), Inches(0.3), Inches(1.3), Inches(0.4),
             "NoQx", size=14, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

def add_footer(slide, page_num, total):
    add_text(slide, Inches(0.5), Inches(7.15), Inches(6), Inches(0.3),
             "NoQx Launch Plan · Prepared 2026-05-22", size=9, color=INK_3)
    add_text(slide, Inches(7), Inches(7.15), Inches(5.8), Inches(0.3),
             f"{page_num} / {total}", size=9, color=INK_3, align=PP_ALIGN.RIGHT)

def add_table(slide, left, top, width, headers, rows, *, header_fill=PURPLE,
              header_text=WHITE, row_alt=LIGHT_BG, font_size=11, header_size=12):
    cols = len(headers)
    n_rows = len(rows) + 1
    row_h = Inches(0.36)
    tbl_shape = slide.shapes.add_table(n_rows, cols, left, top, width, row_h * n_rows)
    tbl = tbl_shape.table
    # Set column widths proportionally; allow caller to tune by editing afterward
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = str(h)
        r.font.size = Pt(header_size); r.font.bold = True
        r.font.color.rgb = header_text
    for ri, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(ri + 1, c)
            if ri % 2 == 1:
                cell.fill.solid(); cell.fill.fore_color.rgb = row_alt
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.05)
            tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(font_size); r.font.color.rgb = INK_DARK
    return tbl

# ── Slide 1: Title ────────────────────────────────────────────────────────
def slide_title():
    s = prs.slides.add_slide(BLANK)
    add_solid_bg(s, PURPLE_DARK)
    if LOGO.exists():
        s.shapes.add_picture(str(LOGO), Inches(5.7), Inches(1.6), Inches(2), Inches(2))
    add_text(s, Inches(0), Inches(3.8), Inches(13.33), Inches(0.9),
             "NoQx Launch Plan", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(4.8), Inches(13.33), Inches(0.5),
             "Skip the queue. Pre-order. Pickup.", size=18, color=RGBColor(0xc4, 0xb5, 0xfd),
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(6.2), Inches(13.33), Inches(0.4),
             "Path from code-ready to live in stores · 2026-05-22",
             size=14, color=RGBColor(0x94, 0xa3, 0xb8), align=PP_ALIGN.CENTER)

# ── Slide 2: Where we are today ──────────────────────────────────────────
def slide_state():
    s = prs.slides.add_slide(BLANK)
    add_solid_bg(s, WHITE)
    add_header_strip(s, "Where we are today", "Code-ready, infrastructure live, store accounts pending")
    add_text(s, Inches(0.5), Inches(1.2), Inches(6), Inches(0.4),
             "DONE", size=14, bold=True, color=GREEN)
    done = [
        "Production web app live at https://noqx.co.in",
        "Cloudflare proxy + SSL Full mode",
        "Supabase production schema + auth + RLS verified",
        "Both Android APKs build, sign, and run on device",
        "iOS Capacitor scaffold compile-checks pass",
        "All legal pages live (Privacy, Terms, Refund, Shipping, Contact)",
        "Email infrastructure operational (Resend + Cloudflare routing)",
        "Brand assets: logo, splash, all icon sizes generated",
        "Store-listing copy + data-safety answers prepared",
        "19/19 production smoke tests passing",
        "Razorpay code integration complete (test mode active)",
        "Mobile keystores generated + backed up (OneDrive)",
    ]
    y = Inches(1.55)
    for item in done:
        add_text(s, Inches(0.7), y, Inches(6), Inches(0.28),
                 f"✓  {item}", size=11, color=INK_DARK)
        y += Inches(0.3)
    add_text(s, Inches(7), Inches(1.2), Inches(6), Inches(0.4),
             "PENDING (operator action)", size=14, bold=True, color=AMBER)
    pending = [
        "Apple Developer Program enrolment ($99/yr)",
        "Google Play Console signup ($25 one-time)",
        "Razorpay KYC + live keys deployment",
        "Supabase Pro upgrade ($25/mo) for DB backups",
        "App store listings created + screenshots captured",
        "Real canteen seeded (placeholder still live)",
        "GST registration decision + flag set",
        "FSSAI license number displayed on canteen pages",
        "Production APK flipped + uploaded to Internal track",
        "iOS signing material (after Apple Dev approves)",
        "TestFlight build + verification",
        "Submit for App Review (Play + Apple)",
    ]
    y = Inches(1.55)
    for item in pending:
        add_text(s, Inches(7.2), y, Inches(6), Inches(0.28),
                 f"○  {item}", size=11, color=INK_2)
        y += Inches(0.3)
    add_footer(s, 2, 10)

# ── Slide 3: 14-day timeline ─────────────────────────────────────────────
def slide_timeline():
    s = prs.slides.add_slide(BLANK)
    add_solid_bg(s, WHITE)
    add_header_strip(s, "14-day timeline", "From sign-ups to live in both stores")
    rows = [
        ["Day 0 (Today)", "Open 4 accounts in parallel", "Apple Dev + Play + Razorpay KYC + Supabase Pro", "$124 + $25/mo"],
        ["Day 1-2", "Wait", "Apple Dev review (24-48h)", "$0"],
        ["Day 2-3", "iOS signing setup", "Cert + Profile + ASC API key → 6 GitHub secrets → TestFlight build", "$0"],
        ["Day 3-7", "Wait + setup", "Razorpay KYC reviewed (3-7d). Meanwhile build Play listings.", "$0"],
        ["Day 5-7", "Live keys + screenshots", "Razorpay live deployed, ₹1 test order. Capture 6+4 screenshots.", "₹1 in fees"],
        ["Day 7", "Production flip", "Real canteen seeded. Capacitor flip to noqx.co.in. Production APKs built.", "$0"],
        ["Day 7-10", "Internal testing", "Upload AAB to Play Internal, IPA to TestFlight. Invite operator + client.", "$0"],
        ["Day 10-12", "Submit for review", "Play review (1-3d). Apple review (24-48h).", "$0"],
        ["Day 12-14", "Approved → live", "Rollout 100% on both stores. Open the canteen.", "\U0001f7e2 LIVE"],
    ]
    tbl = add_table(s, Inches(0.5), Inches(1.4), Inches(12.3),
                    ["When", "What", "Detail", "Cost"], rows, font_size=10, header_size=11)
    # Tune column widths
    tbl.columns[0].width = Inches(1.4)
    tbl.columns[1].width = Inches(2.2)
    tbl.columns[2].width = Inches(7.0)
    tbl.columns[3].width = Inches(1.7)
    add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
             "Critical path is external review wait: Apple Dev 24-48h, Razorpay KYC 3-7 days, app store reviews 1-3 days each.",
             size=11, color=INK_3)
    add_footer(s, 3, 10)

# ── Slide 4: Costs (upfront) ─────────────────────────────────────────────
def slide_costs_upfront():
    s = prs.slides.add_slide(BLANK)
    add_solid_bg(s, WHITE)
    add_header_strip(s, "Upfront costs (one-time)", "Pay these once to get to launch")
    rows = [
        ["Apple Developer Program", "$99/yr", "Required for App Store + TestFlight"],
        ["Google Play Console", "$25 once", "Required to publish on Play Store"],
        ["Razorpay account", "FREE", "2% transaction fee charged per payment"],
        ["Supabase Pro (first month)", "$25", "Database backups + larger DB quota"],
        ["Resend Pro (first month)", "$20", "50k emails / month included"],
        ["Railway Pro (already paying)", "$20", "Hosting infrastructure"],
        ["", "", ""],
        ["TOTAL DAY-1 OUT-OF-POCKET", "$214", "$124 one-time + $90 first month subscriptions"],
    ]
    tbl = add_table(s, Inches(1.5), Inches(1.5), Inches(10.3),
                    ["Service", "Cost", "Why"], rows, font_size=12, header_size=13)
    tbl.columns[0].width = Inches(3.5)
    tbl.columns[1].width = Inches(1.8)
    tbl.columns[2].width = Inches(5.0)
    # Highlight the total row
    for c in range(3):
        cell = tbl.cell(8, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = PURPLE
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = WHITE
                r.font.bold = True
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
             "After day 1, you've spent $214. Apple + Play don't repeat; Supabase + Resend + Railway are monthly.",
             size=11, color=INK_3)
    add_footer(s, 4, 10)

# ── Slide 5: Costs (monthly recurring at scale) ──────────────────────────
def slide_costs_monthly():
    s = prs.slides.add_slide(BLANK)
    add_solid_bg(s, WHITE)
    add_header_strip(s, "Monthly costs at scale",
                     "What you'll pay each month as the app grows")
    rows = [
        ["1,000 DAU", "$20", "$25", "$20", "$0", "$5", "$70/mo"],
        ["5,000 DAU", "$50", "$35", "$99", "$0", "$20", "$200/mo"],
        ["15,000 DAU", "$100", "$85", "$399", "$0", "$50", "$630/mo"],
        ["15k w/ FCM push", "$100", "$85", "$99", "$0", "$50", "$330/mo"],
        ["50,000 DAU", "$300", "$135", "Enterprise", "$20", "$150", "$1,000+/mo"],
    ]
    tbl = add_table(s, Inches(0.5), Inches(1.4), Inches(12.3),
                    ["Daily active students", "Railway", "Supabase", "Resend", "Cloudflare", "SMS", "Total"],
                    rows, font_size=11, header_size=11)
    tbl.columns[0].width = Inches(2.4)
    for c in range(1, 7):
        tbl.columns[c].width = Inches(1.65)
    # Highlight 15k row
    for c in range(7):
        cell = tbl.cell(3, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_BG
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.color.rgb = PURPLE
    add_text(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.5),
             "At 15k DAU: re-enable push notifications (FCM, currently disabled) to drop email volume by ~70% → saves ~$300/mo.",
             size=12, color=INK_DARK)
    add_text(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.5),
             "PLUS Razorpay 2% fee on every order. At 15k DAU × ₹80 avg order × 30 days = ₹36L gross/mo → ₹72k Razorpay fee.",
             size=12, color=INK_DARK)
    add_text(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5),
             "All services are usage-based: you don't pay for capacity you don't use. Scale up only as you cross thresholds.",
             size=11, color=INK_3)
    add_footer(s, 5, 10)

# ── Slide 6: What we need from the client ────────────────────────────────
def slide_from_client():
    s = prs.slides.add_slide(BLANK)
    add_solid_bg(s, WHITE)
    add_header_strip(s, "What we need from the client", "Decisions + content to unblock submission")
    add_text(s, Inches(0.5), Inches(1.2), Inches(6), Inches(0.4),
             "DECISIONS", size=14, bold=True, color=PURPLE)
    decisions = [
        "Apple Developer enrolment type: Individual or Organization (DUNS required)",
        "Razorpay account: KYC documents ready (PAN, GSTIN, bank, director ID)",
        "GST registration status: registered yet? If yes, share GSTIN",
        "Cookie consent: not needed (DECIDED, see privacy policy §8)",
        "FSSAI license number for canteen pages",
        "Real first partner canteen: name, address, college, opening hours",
        "Payment terms with canteen partner (revenue share %)",
    ]
    y = Inches(1.55)
    for d in decisions:
        add_text(s, Inches(0.7), y, Inches(6.0), Inches(0.4),
                 f"•  {d}", size=11, color=INK_DARK)
        y += Inches(0.4)
    add_text(s, Inches(7), Inches(1.2), Inches(6), Inches(0.4),
             "CONTENT TO PROVIDE", size=14, bold=True, color=PURPLE)
    content = [
        "Real first canteen menu (items + prices + categories)",
        "Real canteen photos (3-5 per canteen for student-facing display)",
        "App Store \"What's New\" copy for v1.0",
        "Demo student credentials for App Store reviewer",
        "Privacy contact email + grievance officer name (DPDPA §10)",
        "Apple Developer account credentials (for signing material upload)",
        "Play Console account credentials (for AAB upload)",
        "Razorpay live API keys (after KYC approves)",
    ]
    y = Inches(1.55)
    for c in content:
        add_text(s, Inches(7.2), y, Inches(6.0), Inches(0.4),
                 f"•  {c}", size=11, color=INK_DARK)
        y += Inches(0.4)
    add_footer(s, 6, 10)

# ── Slide 7: What's already done on dev side ─────────────────────────────
def slide_dev_complete():
    s = prs.slides.add_slide(BLANK)
    add_solid_bg(s, WHITE)
    add_header_strip(s, "What's complete on the engineering side",
                     "All code paths tested, documented, and ready")
    items = [
        ("Web application", "Live at noqx.co.in. 19/19 smoke tests pass. Cloudflare proxy active."),
        ("Mobile builds", "Both Android APKs build, sign, install. iOS scaffold compile-checks."),
        ("Brand assets", "NoQx logo (with glittering) regenerated across 38 sizes. Splash + feature graphic."),
        ("Native role isolation", "Student app hides canteen-login UI. Worker app hides student links."),
        ("Email delivery", "Resend wired + verified. Password reset OTP confirmed live."),
        ("Email routing", "Cloudflare → support/grievance/privacy@noqx.co.in → forwards verified."),
        ("Payment integration", "Razorpay test mode working. Live key flip is a 60-second change."),
        ("Data safety / privacy", "Privacy policy DPDPA-compliant. No third-party trackers. Audited."),
        ("Operational runbooks", "On-call, rollback, launch-day, FCM re-enable — all written."),
        ("Smoke + verify scripts", "smoke-test-prod (19 checks). verify-razorpay-flow (5 checks). launch-day-stats."),
        ("Deep linking", "Android App Links: noqx.co.in/dashboard/* opens directly in app."),
        ("UI fixes verified", "Worker header safe-area. Old-order invoice fallback. Worker logout routing."),
    ]
    y = Inches(1.2)
    for title, desc in items:
        add_text(s, Inches(0.5), y, Inches(3.5), Inches(0.3),
                 f"✓  {title}", size=12, bold=True, color=GREEN)
        add_text(s, Inches(4.2), y, Inches(8.6), Inches(0.3),
                 desc, size=11, color=INK_DARK)
        y += Inches(0.45)
    add_footer(s, 7, 10)

# ── Slide 8: Risks + mitigation ──────────────────────────────────────────
def slide_risks():
    s = prs.slides.add_slide(BLANK)
    add_solid_bg(s, WHITE)
    add_header_strip(s, "Risks + mitigation", "What could delay launch and how we handle it")
    rows = [
        ["Apple Dev review stuck", "Medium", "Apple manual verification rarely takes >48h. Follow-up via dev portal support if delayed."],
        ["Razorpay KYC rejection", "Low", "Documents pre-checked. If rejected, re-submit with clearer scans — 24h re-review."],
        ["App Store rejection round 1", "Medium", "Common: demo credentials missing, IAP confusion. Resubmit fixes in 24-48h."],
        ["Play Store rejection round 1", "Low", "Data-safety form pre-prepared. Mostly accepted on first try."],
        ["Razorpay live test fails", "Low", "Test mode flow validated; live flow uses same code path. Rollback to test mode in 60 sec if needed."],
        ["Cloudflare DNS propagation", "Already done", "Zone live since 2026-05-20. No further action."],
        ["DB data loss before backups", "HIGH → mitigated", "Supabase Pro upgrade BLOCKS launch. Day-0 action item."],
        ["First-week traffic spike", "Low (1 canteen)", "Cloudflare absorbs at edge. Railway scales to 70% CPU before we'd notice."],
        ["FCM push not enabled", "Accepted v1", "Users refresh app to see updates. Re-enable in v1.1 per FCM_REENABLEMENT.md."],
    ]
    tbl = add_table(s, Inches(0.5), Inches(1.4), Inches(12.3),
                    ["Risk", "Likelihood", "Mitigation"], rows, font_size=10, header_size=11)
    tbl.columns[0].width = Inches(3.5)
    tbl.columns[1].width = Inches(1.8)
    tbl.columns[2].width = Inches(7.0)
    add_footer(s, 8, 10)

# ── Slide 9: Post-launch monitoring ──────────────────────────────────────
def slide_post_launch():
    s = prs.slides.add_slide(BLANK)
    add_solid_bg(s, WHITE)
    add_header_strip(s, "Post-launch monitoring", "First 30 days are critical")
    add_text(s, Inches(0.5), Inches(1.3), Inches(6.2), Inches(0.4),
             "DAILY (first 7 days)", size=14, bold=True, color=PURPLE)
    daily = [
        "Run launch-day-stats script (orders, payments, signups)",
        "Check Razorpay dashboard for failed payments (target <5%)",
        "Check Resend dashboard for email bounce rate (target <2%)",
        "Review UptimeRobot for any downtime alerts",
        "Triage any customer support tickets <4h response",
    ]
    y = Inches(1.7)
    for item in daily:
        add_text(s, Inches(0.7), y, Inches(6), Inches(0.4),
                 f"•  {item}", size=11, color=INK_DARK)
        y += Inches(0.4)
    add_text(s, Inches(7), Inches(1.3), Inches(6.2), Inches(0.4),
             "WEEKLY", size=14, bold=True, color=PURPLE)
    weekly = [
        "Review Railway usage (CPU, memory, egress trending)",
        "Review Supabase DB size + query performance",
        "Customer retention metric (DAU, WAU)",
        "Operational retro: any new edge cases? Update runbooks",
        "Plan v1.1 fixes based on production feedback",
    ]
    y = Inches(1.7)
    for item in weekly:
        add_text(s, Inches(7.2), y, Inches(6), Inches(0.4),
                 f"•  {item}", size=11, color=INK_DARK)
        y += Inches(0.4)
    add_text(s, Inches(0.5), Inches(4.6), Inches(12.3), Inches(0.4),
             "FIRST 30 DAYS: enable FCM push notifications (FCM_REENABLEMENT.md)",
             size=13, bold=True, color=AMBER)
    add_text(s, Inches(0.5), Inches(5.05), Inches(12.3), Inches(0.4),
             "Cuts Resend email volume by ~70% (saves ~$300/mo at 15k DAU). Requires Firebase project + 7 implementation steps.",
             size=11, color=INK_2)
    add_text(s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.4),
             "QUARTERLY",
             size=14, bold=True, color=PURPLE)
    add_text(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.4),
             "•  Audit Cloudflare WAF rules · rotate keystore backup verification · review pricing-tier upgrade thresholds",
             size=11, color=INK_DARK)
    add_footer(s, 9, 10)

# ── Slide 10: Next 24 hours ──────────────────────────────────────────────
def slide_next_24h():
    s = prs.slides.add_slide(BLANK)
    add_solid_bg(s, PURPLE_DARK)
    add_text(s, Inches(0), Inches(0.7), Inches(13.33), Inches(0.7),
             "Next 24 hours", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(1.5), Inches(13.33), Inches(0.5),
             "These four sign-ups start all wait-clocks in parallel. Do them this morning.",
             size=16, color=RGBColor(0xc4, 0xb5, 0xfd), align=PP_ALIGN.CENTER)
    cards = [
        ("1. Apple Developer", "$99/yr", "developer.apple.com/programs/enroll", "24-48h activation"),
        ("2. Razorpay KYC", "Free signup", "dashboard.razorpay.com", "3-7 day review"),
        ("3. Play Console", "$25 one-time", "play.google.com/console/signup", "Same-day activation"),
        ("4. Supabase Pro", "$25/mo", "supabase.com/dashboard", "Immediate"),
    ]
    # Two-column grid
    positions = [
        (Inches(0.7), Inches(2.5)),
        (Inches(7.0), Inches(2.5)),
        (Inches(0.7), Inches(4.7)),
        (Inches(7.0), Inches(4.7)),
    ]
    for i, ((title, cost, url, eta), (left, top)) in enumerate(zip(cards, positions)):
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.7), Inches(1.9))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x2d, 0x25, 0x55)
        card.line.fill.background()
        card.adjustments[0] = 0.07
        add_text(s, left + Inches(0.3), top + Inches(0.15), Inches(5.3), Inches(0.5),
                 title, size=20, bold=True, color=WHITE)
        add_text(s, left + Inches(0.3), top + Inches(0.6), Inches(5.3), Inches(0.4),
                 cost, size=14, color=RGBColor(0xc4, 0xb5, 0xfd))
        add_text(s, left + Inches(0.3), top + Inches(1.0), Inches(5.3), Inches(0.4),
                 url, size=12, color=RGBColor(0xa7, 0xa3, 0xff))
        add_text(s, left + Inches(0.3), top + Inches(1.4), Inches(5.3), Inches(0.4),
                 f"⏱  {eta}", size=12, color=RGBColor(0x94, 0xa3, 0xb8))
    add_text(s, Inches(0), Inches(6.8), Inches(13.33), Inches(0.4),
             "Total this morning: $124 + 90 minutes of clicking. After that you're waiting on external reviews.",
             size=13, color=RGBColor(0xc4, 0xb5, 0xfd), align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(7.15), Inches(13.33), Inches(0.3),
             "Detailed steps in docs/LAUNCH_CRITICAL_PATH.md",
             size=10, color=RGBColor(0x94, 0xa3, 0xb8), align=PP_ALIGN.CENTER)

# Build deck
slide_title()
slide_state()
slide_timeline()
slide_costs_upfront()
slide_costs_monthly()
slide_from_client()
slide_dev_complete()
slide_risks()
slide_post_launch()
slide_next_24h()

prs.save(str(OUT))
print(f"wrote {OUT} ({OUT.stat().st_size // 1024} KB, {len(prs.slides)} slides)")
